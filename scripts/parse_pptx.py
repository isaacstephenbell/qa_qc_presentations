import sys
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_from_shape(shape):
    """Recursively extracts text from a shape and its sub-shapes."""
    texts = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            texts.extend(extract_text_from_shape(sub_shape))

    elif shape.has_text_frame:
        if shape.text_frame.text.strip():
            texts.append(shape.text_frame.text.strip())

    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame and cell.text_frame.text.strip():
                    texts.append(cell.text_frame.text.strip())

    elif hasattr(shape, 'chart') and shape.chart:
        chart = shape.chart
        if chart.has_title and chart.chart_title.has_text_frame:
            if chart.chart_title.text_frame.text.strip():
                texts.append(chart.chart_title.text_frame.text.strip())

    elif hasattr(shape, "text") and shape.text and shape.text.strip():
        texts.append(shape.text.strip())

    return texts

def parse_pptx(file_path):
    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        
        all_texts = []
        for shape in slide.shapes:
            all_texts.extend(extract_text_from_shape(shape))
        
        # Remove duplicates while preserving order and filter out the title
        seen = {title} if title else set()
        unique_texts = []
        for text in all_texts:
            if text and text not in seen:
                seen.add(text)
                unique_texts.append(text)

        bullets = [{"text": text, "is_bold": False, "indent_level": 0} for text in unique_texts]

        slides.append({
            "slide_number": i + 1,
            "title": title,
            "bullets": bullets
        })

    print(json.dumps(slides, indent=2))

if __name__ == "__main__":
    try:
        if len(sys.argv) < 2:
            raise ValueError("Missing file path argument.")
        parse_pptx(sys.argv[1])
    except Exception as e:
        print(json.dumps({ "error": str(e) }), file=sys.stderr)
        sys.exit(1) 